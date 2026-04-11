import os
from tools.registry import tool

@tool(
    name="create_ppt",
    description="Creates a new PowerPoint presentation (.pptx) with the specified slides.",
    parameters={
        "type": "object",
        "properties": {
            "file_path": {"type": "string", "description": "The path where the .pptx file will be saved."},
            "slides_content": {
                "type": "array",
                "description": "A list of dictionaries representing slides. Each dict should have a 'title' string and an optional 'content' list of strings (for bullet points).",
                "items": {
                    "type": "object",
                    "properties": {
                        "title": {"type": "string"},
                        "content": {"type": "array", "items": {"type": "string"}}
                    },
                    "required": ["title"]
                }
            }
        },
        "required": ["file_path", "slides_content"]
    },
    requires_permission=True
)
def create_ppt(file_path: str, slides_content: list) -> str:
    """Creates a PPTX with the given slides."""
    try:
        from pptx import Presentation
        prs = Presentation()
        
        # Slide layout 1 usually corresponds to a Title and Content slide
        title_content_layout = prs.slide_layouts[1]
        # Slide layout 0 usually corresponds to a Title Slide
        title_slide_layout = prs.slide_layouts[0]
        
        for i, slide_data in enumerate(slides_content):
            title = slide_data.get("title", f"Slide {i+1}")
            bullets = slide_data.get("content", [])
            
            # If it's the first slide and no bullets, make it a title slide
            if i == 0 and not bullets:
                slide = prs.slides.add_slide(title_slide_layout)
                slide.shapes.title.text = title
            else:
                slide = prs.slides.add_slide(title_content_layout)
                slide.shapes.title.text = title
                
                if bullets:
                    tf = slide.placeholders[1].text_frame
                    for j, bullet in enumerate(bullets):
                        if j == 0:
                            tf.text = bullet
                        else:
                            p = tf.add_paragraph()
                            p.text = bullet
                            
        prs.save(file_path)
        return f"Successfully created presentation at '{file_path}'."
    except ImportError:
        return "Error: 'python-pptx' package is not installed. Please run: pip install python-pptx"
    except Exception as e:
        return f"Error creating PowerPoint presentation: {e}"
